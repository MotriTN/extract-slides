
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import img2pdf

def convert_images_to_pdf_pptx(input_folder, pdf_output, pptx_output):
    """
    Convert JPG images in a folder to PDF and PPTX files.
    
    Args:
        input_folder (str): Path to folder containing JPG images
        pdf_output (str): Output PDF filename
        pptx_output (str): Output PPTX filename
    """
    
    # Get all JPG files from input folder, sorted by name
    image_files = sorted([f for f in os.listdir(input_folder) if f.lower().endswith('.jpg')])
    
    if not image_files:
        print(f"No JPG images found in {input_folder}")
        return
    
    print(f"Found {len(image_files)} images to process...")
    
    # Convert to PDF
    print("Creating PDF...")
    with open(pdf_output, "wb") as f:
        # Create list of image paths with full paths
        img_paths = [os.path.join(input_folder, img) for img in image_files]
        # Convert to PDF
        f.write(img2pdf.convert(img_paths))
    
    # Convert to PPTX
    print("Creating PowerPoint...")
    prs = Presentation()
    
    for img_file in image_files:
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]  # 6 is the layout code for a blank slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add picture to slide
        img_path = os.path.join(input_folder, img_file)
        left = top = Inches(0)  # Position at top-left corner
        pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
    
    # Save the PowerPoint
    prs.save(pptx_output)
    
    print(f"Conversion complete! PDF saved as {pdf_output}, PowerPoint saved as {pptx_output}")

if __name__ == "__main__":
    # Input and output file paths
    input_folder = "unique_slides"  # Folder containing JPG images
    pdf_output = "unique_slides.pdf"       # Output PDF filename
    pptx_output = "unique_slides.pptx"     # Output PPTX filename
    
    # Check if input folder exists
    if not os.path.exists(input_folder):
        print(f"Error: Input folder '{input_folder}' does not exist.")
    else:
        convert_images_to_pdf_pptx(input_folder, pdf_output, pptx_output)
